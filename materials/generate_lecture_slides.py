#!/usr/bin/env python3
"""중학교 선생님 AI 실습 연수 — 콘텐츠 풍부 PPT + PDF (밝은 배경, 큰 글씨, 이미지)"""

from __future__ import annotations

import textwrap
import urllib.request
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

OUT_DIR = Path(__file__).parent
IMG_DIR = OUT_DIR / "images"
PPTX_PATH = OUT_DIR / "AI_lecture_script.pptx"
PDF_PATH = OUT_DIR / "AI_lecture_script.pdf"
FONT_PATH = OUT_DIR / "fonts" / "NotoSansKR.ttf"
FONT_URL = "https://github.com/google/fonts/raw/main/ofl/notosanskr/NotoSansKR%5Bwght%5D.ttf"

# 색상
C = {
    "blue": (59, 130, 246),
    "blue_d": (30, 64, 175),
    "blue_l": (219, 234, 254),
    "pink": (244, 114, 182),
    "pink_l": (252, 231, 243),
    "slate": (51, 65, 85),
    "slate_l": (100, 116, 139),
    "white": (255, 255, 255),
    "bg": (248, 250, 252),
    "bg_blue": (240, 249, 255),
    "emerald": (5, 150, 105),
    "emerald_l": (209, 250, 229),
    "amber": (217, 119, 6),
    "amber_l": (254, 243, 199),
    "prompt_bg": (30, 41, 59),
    "orange": (234, 88, 12),
}

RGB = {k: RGBColor(*v) for k, v in C.items()}

# ── 폰트 ──────────────────────────────────────────────
def ensure_font():
    if FONT_PATH.exists():
        return
    FONT_PATH.parent.mkdir(parents=True, exist_ok=True)
    print("Downloading Noto Sans KR...")
    urllib.request.urlretrieve(FONT_URL, FONT_PATH)


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    ensure_font()
    return ImageFont.truetype(str(FONT_PATH), size)


# ── 이미지 생성 ───────────────────────────────────────
def _rounded_rect(draw, xy, radius, fill):
    x0, y0, x1, y1 = xy
    draw.rounded_rectangle(xy, radius=radius, fill=fill)


def gen_roadmap(path: Path):
    w, h = 1200, 400
    img = Image.new("RGB", (w, h), C["bg_blue"])
    draw = ImageDraw.Draw(img)
    f = pil_font(22)
    fb = pil_font(26, True)
    steps = [
        ("0~5분", "아이온\n접속", C["blue"]),
        ("5~20분", "NotebookLM\n실습", C["blue"]),
        ("20~32분", "부서/담임\n업무", C["pink"]),
        ("32~45분", "Gemini\n심화", C["pink"]),
        ("45~48분", "워크스페이스+", C["emerald"]),
        ("48~50분", "마무리", C["emerald"]),
    ]
    bw = 170
    for i, (time, label, col) in enumerate(steps):
        x = 30 + i * (bw + 15)
        _rounded_rect(draw, (x, 80, x + bw, 320), 20, col)
        draw.text((x + bw // 2, 110), time, fill=C["white"], font=fb, anchor="mm")
        for j, line in enumerate(label.split("\n")):
            draw.text((x + bw // 2, 180 + j * 32), line, fill=C["white"], font=f, anchor="mm")
        if i < len(steps) - 1:
            draw.polygon([(x + bw + 3, 200), (x + bw + 18, 190), (x + bw + 18, 210)], fill=col)
    draw.text((w // 2, 30), "50분 강의 로드맵", fill=C["blue_d"], font=pil_font(32, True), anchor="mm")
    img.save(path)


def gen_quick_start(path: Path):
    w, h = 1100, 380
    img = Image.new("RGB", (w, h), C["white"])
    draw = ImageDraw.Draw(img)
    steps = [
        ("1", "광주아이온 접속", "aion.gen.go.kr\nteacher@genedu.kr", C["blue"], C["blue_l"]),
        ("2", "프롬프트 복사", "[복사] 버튼 →\nAI에 붙여넣기", C["pink"], C["pink_l"]),
        ("3", "확인 후 수정", "AI 결과 = 초안\n선생님이 완성", C["emerald"], C["emerald_l"]),
    ]
    for i, (num, title, desc, col, bg) in enumerate(steps):
        x = 30 + i * 360
        _rounded_rect(draw, (x, 40, x + 330, 340), 24, bg)
        draw.ellipse((x + 20, 60, x + 80, 120), fill=col)
        draw.text((x + 50, 90), num, fill=C["white"], font=pil_font(30, True), anchor="mm")
        draw.text((x + 100, 75), title, fill=col, font=pil_font(24, True))
        for j, line in enumerate(desc.split("\n")):
            draw.text((x + 30, 150 + j * 36), line, fill=C["slate"], font=pil_font(20))
    img.save(path)


def gen_aion_flow(path: Path):
    w, h = 1100, 420
    img = Image.new("RGB", (w, h), C["bg"])
    draw = ImageDraw.Draw(img)
    draw.text((w // 2, 25), "광주아이온 통합 인증(SSO) 흐름", fill=C["blue_d"], font=pil_font(28, True), anchor="mm")
    boxes = [
        ("teacher@\ngenedu.kr", C["blue_l"], C["blue_d"]),
        ("광주아이온\n대시보드", C["blue"], C["white"]),
        ("Gemini", C["pink_l"], C["pink"]),
        ("NotebookLM", C["pink_l"], C["pink"]),
        ("클래스룸·드라이브", C["emerald_l"], C["emerald"]),
    ]
    bw, bh = 190, 100
    y = 120
    for i, (text, bg, fg) in enumerate(boxes):
        x = 40 + i * 210
        _rounded_rect(draw, (x, y, x + bw, y + bh), 16, bg)
        for j, line in enumerate(text.split("\n")):
            draw.text((x + bw // 2, y + 35 + j * 28), line, fill=fg, font=pil_font(18, True), anchor="mm")
        if i < len(boxes) - 1:
            draw.text((x + bw + 8, y + bh // 2), "→", fill=C["slate_l"], font=pil_font(28, True), anchor="lm")
    features = ["나의 강의실", "학사일정", "교육 자료실", "학습 분석 대시보드"]
    for i, feat in enumerate(features):
        x = 80 + i * 250
        _rounded_rect(draw, (x, 280, x + 220, 360), 12, C["white"])
        draw.text((x + 110, 320), feat, fill=C["slate"], font=pil_font(18, True), anchor="mm")
    img.save(path)


def gen_notebooklm_flow(path: Path):
    w, h = 1000, 350
    img = Image.new("RGB", (w, h), C["bg_blue"])
    draw = ImageDraw.Draw(img)
    draw.text((w // 2, 28), "NotebookLM 따라하기", fill=C["blue_d"], font=pil_font(28, True), anchor="mm")
    steps = ["+ 새 노트", "+ 소스 추가\n(PDF·HWP)", "질문하기", "답변 확인"]
    for i, s in enumerate(steps):
        x = 50 + i * 230
        _rounded_rect(draw, (x, 90, x + 190, 200), 18, C["white"])
        for j, line in enumerate(s.split("\n")):
            draw.text((x + 95, 130 + j * 30), line, fill=C["blue_d"], font=pil_font(20, True), anchor="mm")
        if i < 3:
            draw.text((x + 200, 145), "→", fill=C["blue"], font=pil_font(32, True))
    uses = ["학교 규정·일정", "세특 초안", "평가계획", "문항·루브릭"]
    for i, u in enumerate(uses):
        x = 70 + i * 230
        _rounded_rect(draw, (x, 240, x + 190, 310), 12, C["blue"])
        draw.text((x + 95, 275), u, fill=C["white"], font=pil_font(17, True), anchor="mm")
    img.save(path)


def gen_tool_choice(path: Path):
    w, h = 1000, 380
    img = Image.new("RGB", (w, h), C["white"])
    draw = ImageDraw.Draw(img)
    draw.text((w // 2, 30), "30초 도구 선택법", fill=C["blue_d"], font=pil_font(30, True), anchor="mm")
    draw.text((w // 2, 80), "참고할 학교 문서·파일이 있나요?", fill=C["slate"], font=pil_font(22, True), anchor="mm")
    _rounded_rect(draw, (80, 130, 430, 330), 20, C["blue_l"])
    draw.text((255, 170), "있으면", fill=C["blue_d"], font=pil_font(24, True), anchor="mm")
    draw.text((255, 220), "NotebookLM", fill=C["blue"], font=pil_font(32, True), anchor="mm")
    draw.text((255, 280), "규정·일정·보고서\n기반 답변", fill=C["slate"], font=pil_font(18), anchor="mm")
    draw.text((500, 230), "VS", fill=C["pink"], font=pil_font(36, True), anchor="mm")
    _rounded_rect(draw, (570, 130, 920, 330), 20, C["pink_l"])
    draw.text((745, 170), "없으면", fill=C["blue_d"], font=pil_font(24, True), anchor="mm")
    draw.text((745, 220), "Gemini", fill=C["pink"], font=pil_font(32, True), anchor="mm")
    draw.text((745, 280), "글쓰기·자료 제작\n아이디어 생성", fill=C["slate"], font=pil_font(18), anchor="mm")
    img.save(path)


def gen_slides_7step(path: Path):
    w, h = 1100, 500
    img = Image.new("RGB", (w, h), C["bg"])
    draw = ImageDraw.Draw(img)
    draw.text((w // 2, 22), "수업 자료 제작 7단계", fill=C["blue_d"], font=pil_font(28, True), anchor="mm")
    steps = [
        "1.Gemini\n내용 구성", "2.슬라이드\n제미나이 열기", "3.이미지\n생성(영어)",
        "4.배경\n삽입", "5.Create\nslide", "6.테마\n적용", "7.발표\n스크립트",
    ]
    for i, s in enumerate(steps):
        row, col = i // 4, i % 4
        if i >= 4:
            col = i - 4
            row = 1
        x = 40 + col * 260
        y = 70 + row * 200
        _rounded_rect(draw, (x, y, x + 230, y + 150), 16, C["blue"] if i % 2 == 0 else C["pink"])
        for j, line in enumerate(s.split("\n")):
            draw.text((x + 115, y + 55 + j * 30), line, fill=C["white"], font=pil_font(18, True), anchor="mm")
    img.save(path)


def gen_situation_grid(path: Path):
    w, h = 1100, 450
    img = Image.new("RGB", (w, h), C["bg_blue"])
    draw = ImageDraw.Draw(img)
    draw.text((w // 2, 22), "학사 시기별 활용", fill=C["blue_d"], font=pil_font(28, True), anchor="mm")
    rows = [
        ("3월 학기 초", "학부모 안내·오리엔테이션", "Gemini"),
        ("4~6월", "평가·현장체험학습", "NotebookLM"),
        ("7월", "세특·생기부 기록", "NotebookLM"),
        ("방학", "수업 설계·교육과정", "Gemini"),
        ("연중", "규정·일정·공문 확인", "NotebookLM"),
    ]
    for i, (period, task, tool) in enumerate(rows):
        y = 65 + i * 72
        _rounded_rect(draw, (40, y, w - 40, y + 60), 12, C["white"])
        draw.text((60, y + 30), period, fill=C["blue_d"], font=pil_font(19, True), anchor="lm")
        draw.text((280, y + 30), task, fill=C["slate"], font=pil_font(18), anchor="lm")
        col = C["blue"] if "Notebook" in tool else C["pink"]
        _rounded_rect(draw, (850, y + 12, 1040, y + 48), 8, col)
        draw.text((945, y + 30), tool, fill=C["white"], font=pil_font(16, True), anchor="mm")
    img.save(path)


def generate_all_images():
    IMG_DIR.mkdir(parents=True, exist_ok=True)
    gen_roadmap(IMG_DIR / "roadmap.png")
    gen_quick_start(IMG_DIR / "quick_start.png")
    gen_aion_flow(IMG_DIR / "aion_flow.png")
    gen_notebooklm_flow(IMG_DIR / "notebooklm_flow.png")
    gen_tool_choice(IMG_DIR / "tool_choice.png")
    gen_slides_7step(IMG_DIR / "slides_7step.png")
    gen_situation_grid(IMG_DIR / "situation_grid.png")
    print(f"Images saved to {IMG_DIR}")


# ── 슬라이드 데이터 ───────────────────────────────────
SLIDES: list[dict] = [
    {
        "type": "title",
        "title": "AI와 함께 스마트하게 일하기",
        "subtitle": "중학교 선생님 대상 생성형 AI 실습 연수",
        "lines": ["광주아이온(AI-ON) · Gemini · NotebookLM", "연수 50분  |  실습 가이드 웹페이지 병행"],
        "notes": "안녕하세요. 오늘은 생성형 AI를 수업과 학교 업무에 활용하는 실습 연수입니다. 강의 후에도 따라할 수 있는 가이드 페이지 링크를 드립니다.",
    },
    {
        "type": "image_content",
        "title": "오늘의 목표 & 50분 로드맵",
        "image": "roadmap.png",
        "bullets": [
            "목표 ① 광주아이온 접속 · Gemini · NotebookLM 기본 익히기",
            "목표 ② 수업·업무에 바로 쓸 프롬프트 실습",
            "목표 ③ 강의 후 스스로 활용할 가이드 확보",
        ],
        "notes": "50분 로드맵을 보시면서 오늘 흐름을 안내합니다. 아이온 → NotebookLM → 업무 활용 → Gemini → 마무리 순입니다.",
    },
    {
        "type": "image_content",
        "title": "처음이신가요? 3분 빠른 시작",
        "image": "quick_start.png",
        "bullets": [
            "글자가 작으면 가이드 [가+] 버튼  |  용어는 [AI 기초 도움말]",
            "📌 이 페이지 링크를 저장하면 강의 후에도 언제든 따라할 수 있습니다",
        ],
        "notes": "AI가 낯설어도 3단계만 기억하세요. 로그인 → 복사 → 확인·수정. 지금부터 함께 따라해 봅시다.",
    },
    {
        "type": "split",
        "title": "광주아이온(AI-ON) — 통합 인증 & 접속",
        "time": "0~5분",
        "image": "aion_flow.png",
        "bullets": [
            "아이디: teacher@genedu.kr  →  통합 ID 하나로 전 서비스 연결",
            "메인 화면 위젯 클릭 → Gemini, NotebookLM 이동",
            "[나의 화면 구성하기]로 위젯 배치 조절",
            "나의 강의실 · 학사일정 · 교육 자료실 자동 연동",
        ],
        "notes": "aion.gen.go.kr에 함께 접속합니다. SSO로 한 번 로그인하면 Gemini, NotebookLM, 클래스룸이 연결됩니다.",
    },
    {
        "type": "content",
        "title": "광주아이온 — 학습 분석 대시보드",
        "time": "0~5분",
        "bullets": [
            "활동 분석: 학생 에듀테크 접속 시간·학습 패턴 시각화",
            "퀴즈 리포트: 구글 폼 퀴즈 정답률 분석 → 맞춤 피드백 지원",
            "학습 분포: 주도적 학습 시간 파악 → 수행평가 설계 데이터",
        ],
        "tip": "매뉴얼 24p 참고 — 학생 학습 데이터를 수업 설계에 활용하세요.",
        "notes": "학습 분석 대시보드로 학생들의 학습 패턴을 확인할 수 있습니다. 수행평가 설계 시 참고 자료로 활용하세요.",
    },
    {
        "type": "split",
        "title": "NotebookLM — 나만의 학교 비서",
        "time": "5~20분",
        "image": "notebooklm_flow.png",
        "bullets": [
            "1-1 학교 맞춤 챗봇  |  1-2 세특 초안  |  1-3 평가계획  |  1-4 문항·루브릭",
            "학교 규정·학사일정·공문·보고서를 소스로 업로드",
            "업로드된 자료만 근거로 답변 → 학교 맞춤형!",
        ],
        "notes": "NotebookLM은 선생님이 올린 문서만 기반으로 답합니다. notebooklm.google.com에서 새 노트를 만들어 봅시다.",
    },
    {
        "type": "prompt",
        "title": "1-1. 우리 학교 맞춤형 챗봇",
        "time": "5~20분",
        "steps": ["NotebookLM → + 새 노트", "+ 소스 추가 (규정·일정)", "하단 대화창에 질문"],
        "prompt": "- 2학년 현장체험학습 신청 마감일이 언제야?\n- 스마트폰 사용 적발 시 징계 규정을 알려줘.\n- 올해 여름방학 시작일과 개학일은?",
        "tip": "새 공문이 생길 때마다 소스를 추가하면 챗봇이 계속 똑똑해집니다.",
        "notes": "학교 규정 PDF를 올리고 프롬프트를 복사해 질문해 보세요. 가이드 페이지 [복사] 버튼을 활용합니다.",
    },
    {
        "type": "prompt",
        "title": "1-2. 학생부 기록(세특) 시간 단축",
        "time": "5~20분",
        "steps": ["학생별 노트 생성 (예: 30101 김민준)", "과제·보고서·활동기록 업로드", "세특 초안 요청"],
        "prompt": "이 학생의 탐구 보고서를 바탕으로, 과학적 탐구 역량과 심화 학습 태도가 드러나는 세특 초안을 1~2문장으로 작성해 줘. 데이터 분석 능력이 돋보이도록.",
        "tip": "AI 초안은 참고 자료입니다. 선생님의 관찰을 더해 완성하세요.",
        "notes": "세특 시즌에 학생별 노트를 만들어 두면 매우 유용합니다. 반드시 선생님 관찰을 더해 주세요.",
    },
    {
        "type": "two_prompt",
        "title": "1-3·1-4. 평가 계획 & 문항·루브릭",
        "time": "5~20분",
        "left_title": "1-3 평가 계획 초안",
        "left_prompt": "업로드된 '성취기준'과 '2024 평가계획'을 바탕으로, 2025학년도 '확률과 통계' 1학기 평가 계획 초안을 표 형식으로 작성해 줘. (지필 1회, 수행 2회)",
        "right_title": "1-4 문항 & 루브릭",
        "right_prompt": "(문항) '통합사회 5단' 내용으로 비판적 사고력 평가용 서술형 문항 3개 만들어 줘.\n(기준표) 2번 문항의 채점 기준표를 '문제 이해도', '논리적 근거', '창의적 대안' 항목으로 표로 만들어 줘.",
        "notes": "성취기준 PDF와 작년 평가계획을 소스로 올리면 새 학년도 초안을 빠르게 만들 수 있습니다.",
    },
    {
        "type": "two_prompt",
        "title": "부서별 업무 맞춤형 AI 활용",
        "time": "20~32분",
        "left_title": "📑 교무기획부",
        "left_prompt": "- 2025학년도 1학기 학사일정 주요 행사를 월별 표로 정리해줘.\n- 최근 3년 교원평가 보고서를 바탕으로 올해 연수 계획 시사점 3가지 제안해줘.",
        "right_title": "🔬 교육연구부",
        "right_prompt": "- 국어과 1학년 수행평가 '현대시 비평문 쓰기' 루브릭을 만들어줘.\n- 학교평가 보고서 초안을 위해 성과와 개선점 중심 개요를 작성해줘.",
        "notes": "부서 업무에 맞는 프롬프트를 골라 바로 시도해 보세요.",
    },
    {
        "type": "two_prompt",
        "title": "담임선생님을 위한 AI 치트키",
        "time": "20~32분",
        "left_title": "💌 학부모 소통",
        "left_prompt": "- 2학기 학부모 상담주간 안내 가정통신문 초안을 친근한 톤으로 작성해줘.\n- 김민준 학생 상담일지를 바탕으로 논의할 점 3가지를 요약해줘.",
        "right_title": "🎨 학급 특색 활동",
        "right_prompt": "- 우리 반 협동심을 높일 장기 프로젝트 아이디어 5가지를 제안해줘.\n- 작년 우수사례집을 바탕으로 우리 반 맞춤 활동 3가지를 추천해줘.",
        "notes": "담임 업무에서 가정통신문과 학급 활동 아이디어에 AI를 활용해 보세요.",
    },
    {
        "type": "prompt",
        "title": "Gemini — @기능 & Gems",
        "time": "32~45분",
        "bullets": [
            "@기능: Gemini 채팅에서 @Google 지도 · @Drive · @YouTube 연동",
            "Gems: 자주 쓰는 AI 역할 저장 — 프롬프트 생성기, 양식 생성기",
        ],
        "prompt": "🗺️ @Google 지도 경주 2박3일 수학여행 유적지 중심 동선을 짜줘.\n📂 @Google Drive '2025 평가계획.pdf'를 요약해줘.\n▶️ @YouTube 셰익스피어 '햄릿' 중학생용 영상을 찾아줘.",
        "notes": "@기능으로 실시간 정보를 연동합니다. Gems는 가이드에 시나리오 링크가 있습니다.",
    },
    {
        "type": "prompt",
        "title": "Gemini — Google Forms 설문 자동화",
        "time": "32~45분",
        "steps": [
            "사전 준비(최초 1회): Google 언어 → English (US)",
            "Workspace Labs 가입 (labs.google.com)",
            "Forms → 'Create form with Gemini' → 영어로 요청",
        ],
        "prompt": "Create a survey for 2nd-grade middle school students about their preferred school trip destination. Include choices (Jeju, Gyeongju, Busan, Gangwon-do) and reasons. Please output all text in Korean.",
        "tip": "한글 결과물: 프롬프트 끝에 Please output all text in Korean. 추가",
        "notes": "Forms Gemini는 최초 1회 영어 설정이 필요합니다. 한글 결과는 프롬프트 끝에 Please output all text in Korean.을 붙이세요.",
    },
    {
        "type": "split",
        "title": "Gemini — 수업 자료 제작 7단계",
        "time": "32~45분",
        "image": "slides_7step.png",
        "prompt": "[Step1] '조선 후기 사회의 변화' 주제로 중학생 대상 슬라이드 5장 구성해줘.\n[Step2] A vibrant market scene in late Joseon Dynasty, digital art, presentation background.\n[Step3] 생성한 슬라이드 바탕으로 쉽고 재미있는 발표 스크립트 작성해줘.",
        "notes": "Gemini로 내용 구성 → 슬라이드 제미나이 패널에서 이미지·슬라이드 생성 → 발표 스크립트까지 7단계입니다.",
    },
    {
        "type": "two_prompt",
        "title": "워크스페이스+ 심화 팁",
        "time": "45~48분",
        "left_title": "📊 스프레드시트 분석",
        "left_prompt": "이 스프레드시트 데이터에서 [국어] 점수가 [80점] 이상인 학생들의 특징을 분석해주고, 보완 조언을 작성해줘.",
        "right_title": "🎓 클래스룸 피드백",
        "right_prompt": "업로드된 학생 과제 PDF를 읽고, [성취기준]에 비추어 잘한 점 2가지와 보완점 1가지를 따뜻하게 피드백해줘.",
        "notes": "함수 없이 자연어로 데이터 분석, 클래스룸에서 맞춤 피드백 초안을 받을 수 있습니다.",
    },
    {
        "type": "prompt",
        "title": "[BONUS] 2022 개정 교육과정 수업 설계",
        "time": "45~48분",
        "prompt": "2022 개정 교육과정 [과목명]의 성취기준을 반영하여, 학생들의 [비판적 사고 역량]을 끌어낼 수 있는 수업을 설계해줘. 매 차시 학생이 스스로 질문을 생성하는 활동을 포함해줘.",
        "tip": "역량 중심의 깊이 있는 학습을 AI와 함께 설계하세요!",
        "notes": "방학에 수업을 설계할 때 과목명과 역량만 바꿔서 활용하세요.",
    },
    {
        "type": "split",
        "title": "상황별 AI 활용 가이드 — 강의 후 참고",
        "image": "tool_choice.png",
        "image2": "situation_grid.png",
        "bullets": [
            "규정·일정·세특·평가 → NotebookLM",
            "가정통신문·슬라이드·설문·아이디어 → Gemini",
            "가이드 [상황별 활용] + 검색창으로 프롬프트 바로 찾기",
        ],
        "notes": "강의 후 가장 많은 질문이 '이럴 때 뭘 쓰나요?'입니다. 이 슬라이드를 북마크해 두세요.",
    },
    {
        "type": "closing",
        "title": "마무리 — AI는 선생님의 파트너",
        "cards": [
            ("🔒 개인정보 보호", "학생 실명·연락처 등 민감 정보는 AI에 입력 금지"),
            ("🔍 사실 확인", "AI 생성 정보는 공신력 있는 자료로 재확인"),
            ("✅ 최종 판단", "모든 교육적 판단의 책임은 선생님에게"),
        ],
        "footer": "AI는 선생님을 대체하지 않고, 아이들과의 정서적 교감에 더 집중하도록 돕습니다.",
        "notes": "개인정보 보호, 사실 확인, 최종 판단 세 가지를 꼭 기억해 주세요. 실습 가이드 링크를 공유합니다. Q&A",
    },
]


# ── PPT 빌더 ──────────────────────────────────────────
def set_bg(slide, prs, color_key="white"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGB[color_key]


def add_bar(slide, prs, y=Inches(0), h=Inches(0.1), color="blue"):
    s = slide.shapes.add_shape(1, Inches(0), y, prs.slide_width, h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGB[color]
    s.line.fill.background()


def tb(slide, l, t, w, h, text="", size=22, bold=False, color="slate", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGB[color]
    p.alignment = align
    return tf


def bullets(tf, items, size=22, color="slate", sp=Pt(6)):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = RGB[color]
        p.space_after = sp
        p.level = 0


def add_prompt_box(slide, l, t, w, h, text, size=15):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGB["prompt_bg"]
    shape.line.color.rgb = RGB["slate_l"]
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = RGB["white"]
    p.font.name = "Courier New"


def slide_header(slide, prs, data):
    time_label = data.get("time", "")
    if time_label:
        badge = slide.shapes.add_shape(1, Inches(0.6), Inches(0.35), Inches(1.4), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGB["blue_l"]
        badge.line.fill.background()
        tb(slide, Inches(0.6), Inches(0.32), Inches(1.4), Inches(0.45),
           f"[{time_label}]", size=16, bold=True, color="blue_d", align=PP_ALIGN.CENTER)
    tb(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.8),
       data["title"], size=34, bold=True, color="blue_d")
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.65), Inches(2.8), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGB["pink"]
    bar.line.fill.background()


def add_image(slide, name, l, t, w):
    path = IMG_DIR / name
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w)


def build_pptx():
    generate_all_images()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, data in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        st = data["type"]

        if st == "title":
            set_bg(slide, prs, "bg_blue")
            add_bar(slide, prs, color="blue")
            add_bar(slide, prs, Inches(7.35), Inches(0.08), "pink")
            tb(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
               data["title"], size=48, bold=True, color="blue_d", align=PP_ALIGN.CENTER)
            tb(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.7),
               data["subtitle"], size=26, color="slate", align=PP_ALIGN.CENTER)
            tf = tb(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(1.5), "", size=22, color="slate_l", align=PP_ALIGN.CENTER)
            bullets(tf, data["lines"], size=24, color="slate")

        elif st == "closing":
            set_bg(slide, prs, "bg_blue")
            add_bar(slide, prs, color="pink")
            slide_header(slide, prs, data)
            for i, (title, desc) in enumerate(data["cards"]):
                x = Inches(0.6 + i * 4.2)
                card = slide.shapes.add_shape(1, x, Inches(2.0), Inches(3.9), Inches(2.8))
                card.fill.solid()
                card.fill.fore_color.rgb = RGB["white"]
                card.line.color.rgb = RGB["blue_l"]
                tb(slide, x + Inches(0.2), Inches(2.2), Inches(3.5), Inches(0.6),
                   title, size=24, bold=True, color="blue_d")
                tb(slide, x + Inches(0.2), Inches(3.0), Inches(3.5), Inches(1.5),
                   desc, size=20, color="slate")
            tb(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8),
               data["footer"], size=22, bold=True, color="blue_d", align=PP_ALIGN.CENTER)
            tb(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
               "실습 가이드 링크 공유  →  Q&A", size=24, color="pink", align=PP_ALIGN.CENTER)

        elif st == "image_content":
            set_bg(slide, prs, "white")
            slide_header(slide, prs, data)
            add_image(slide, data["image"], Inches(0.6), Inches(1.9), Inches(12))
            tf = tb(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), "", size=22)
            bullets(tf, data.get("bullets", []), size=22)

        elif st == "split":
            set_bg(slide, prs, "white")
            slide_header(slide, prs, data)
            add_image(slide, data.get("image", ""), Inches(0.6), Inches(1.85), Inches(6.2))
            tf = tb(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(3.5), "", size=22)
            bullets(tf, data.get("bullets", []), size=22)
            if data.get("prompt"):
                add_prompt_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.8), data["prompt"], 14)
            if data.get("image2"):
                add_image(slide, data["image2"], Inches(0.6), Inches(4.8), Inches(6.2))

        elif st == "content":
            set_bg(slide, prs, "bg")
            slide_header(slide, prs, data)
            tf = tb(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.0), "", size=24)
            bullets(tf, data.get("bullets", []), size=24)
            if data.get("tip"):
                tip = slide.shapes.add_shape(1, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0))
                tip.fill.solid()
                tip.fill.fore_color.rgb = RGB["amber_l"]
                tip.line.fill.background()
                tb(slide, Inches(1.0), Inches(5.6), Inches(11.0), Inches(0.8),
                   f"💡 {data['tip']}", size=20, color="amber")

        elif st == "prompt":
            set_bg(slide, prs, "white")
            slide_header(slide, prs, data)
            y = Inches(1.9)
            if data.get("steps"):
                for i, step in enumerate(data["steps"]):
                    x = Inches(0.6 + i * 4.1)
                    box = slide.shapes.add_shape(1, x, y, Inches(3.8), Inches(0.9))
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGB["blue_l"]
                    box.line.fill.background()
                    tb(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.7),
                       f"{i+1}. {step}", size=18, bold=True, color="blue_d")
                y = Inches(3.0)
            if data.get("bullets"):
                tf = tb(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), "", size=22)
                bullets(tf, data["bullets"], size=22)
                y = Inches(4.0)
            add_prompt_box(slide, Inches(0.6), y, Inches(12.0), Inches(2.0), data.get("prompt", ""), 16)
            if data.get("tip"):
                tb(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
                   f"💡 {data['tip']}", size=20, bold=True, color="amber")

        elif st == "two_prompt":
            set_bg(slide, prs, "bg")
            slide_header(slide, prs, data)
            for i, (title_key, prompt_key, x) in enumerate([
                ("left_title", "left_prompt", Inches(0.6)),
                ("right_title", "right_prompt", Inches(6.8)),
            ]):
                card = slide.shapes.add_shape(1, x, Inches(1.9), Inches(6.0), Inches(5.0))
                card.fill.solid()
                card.fill.fore_color.rgb = RGB["white"]
                card.line.color.rgb = RGB["blue_l"] if i == 0 else RGB["pink"]
                tb(slide, x + Inches(0.2), Inches(2.05), Inches(5.6), Inches(0.5),
                   data[title_key], size=24, bold=True, color="blue_d" if i == 0 else "pink")
                add_prompt_box(slide, x + Inches(0.2), Inches(2.7), Inches(5.6), Inches(3.9),
                               data[prompt_key], 15)

        slide.notes_slide.notes_text_frame.text = data.get("notes", "")
        tb(slide, Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.35),
           str(idx), size=14, color="slate_l", align=PP_ALIGN.RIGHT)

    prs.save(PPTX_PATH)
    print(f"PPTX: {PPTX_PATH} ({len(SLIDES)} slides)")


# ── PDF 빌더 ──────────────────────────────────────────
def wrap_text(c, text, font, size, max_w):
    lines = []
    for paragraph in text.split("\n"):
        line = ""
        for ch in paragraph:
            test = line + ch
            if c.stringWidth(test, font, size) <= max_w:
                line = test
            else:
                if line:
                    lines.append(line)
                line = ch
        if line:
            lines.append(line)
    return lines


def build_pdf():
    ensure_font()
    pdfmetrics.registerFont(TTFont("NotoKR", str(FONT_PATH)))
    pw, ph = landscape((297 * mm, 210 * mm))
    c = canvas.Canvas(str(PDF_PATH), pagesize=(pw, ph))
    generate_all_images()

    for i, data in enumerate(SLIDES, 1):
        # 밝은 배경
        c.setFillColor(colors.HexColor("#F8FAFC"))
        c.rect(0, 0, pw, ph, fill=1, stroke=0)
        c.setFillColor(colors.HexColor("#3B82F6"))
        c.rect(0, ph - 12, pw, 12, fill=1, stroke=0)

        st = data["type"]
        y_title = ph - 55

        if st == "title":
            c.setFillColor(colors.HexColor("#EFF6FF"))
            c.rect(0, 0, pw, ph, fill=1, stroke=0)
            c.setFillColor(colors.HexColor("#1E40AF"))
            c.setFont("NotoKR", 32)
            c.drawCentredString(pw / 2, ph - 100, data["title"])
            c.setFont("NotoKR", 18)
            c.setFillColor(colors.HexColor("#334155"))
            c.drawCentredString(pw / 2, ph - 135, data.get("subtitle", ""))
            c.setFont("NotoKR", 16)
            y = ph - 180
            for line in data.get("lines", []):
                c.drawCentredString(pw / 2, y, line)
                y -= 28

        elif st == "closing":
            c.setFillColor(colors.HexColor("#1E40AF"))
            c.setFont("NotoKR", 24)
            c.drawString(40, y_title, data["title"])
            y = ph - 100
            for title, desc in data.get("cards", []):
                c.setFillColor(colors.white)
                c.roundRect(40, y - 80, 240, 90, 8, fill=1, stroke=0)
                c.setFillColor(colors.HexColor("#1E40AF"))
                c.setFont("NotoKR", 14)
                c.drawString(55, y - 20, title)
                c.setFillColor(colors.HexColor("#334155"))
                c.setFont("NotoKR", 11)
                for j, ln in enumerate(wrap_text(c, desc, "NotoKR", 11, 210)):
                    c.drawString(55, y - 42 - j * 16, ln)
                y -= 100
            c.setFont("NotoKR", 13)
            c.setFillColor(colors.HexColor("#1E40AF"))
            c.drawCentredString(pw / 2, 80, data.get("footer", ""))

        else:
            time_label = data.get("time", "")
            if time_label:
                c.setFillColor(colors.HexColor("#DBEAFE"))
                c.roundRect(35, ph - 48, 65, 18, 4, fill=1, stroke=0)
                c.setFillColor(colors.HexColor("#1E40AF"))
                c.setFont("NotoKR", 10)
                c.drawCentredString(67, ph - 42, f"[{time_label}]")

            c.setFillColor(colors.HexColor("#1E40AF"))
            c.setFont("NotoKR", 22)
            c.drawString(40, y_title, data["title"])
            c.setFillColor(colors.HexColor("#F472B6"))
            c.rect(40, y_title - 14, 70, 3, fill=1, stroke=0)

            # 이미지
            img_name = data.get("image")
            if img_name:
                img_path = IMG_DIR / img_name
                if img_path.exists():
                    c.drawImage(str(img_path), 40, 60, width=360, height=120, preserveAspectRatio=True, mask="auto")

            c.setFillColor(colors.HexColor("#334155"))
            c.setFont("NotoKR", 13)
            y = ph - 95
            for b in data.get("bullets", []):
                for ln in wrap_text(c, f"• {b}", "NotoKR", 13, pw - 80):
                    c.drawString(50, y, ln)
                    y -= 18
                y -= 4

            for step in data.get("steps", []):
                c.drawString(50, y, f"→ {step}")
                y -= 18

            prompt = data.get("prompt") or data.get("left_prompt", "")
            if prompt:
                c.setFillColor(colors.HexColor("#1E293B"))
                c.roundRect(40, 30, pw - 80, 55, 6, fill=1, stroke=0)
                c.setFillColor(colors.white)
                c.setFont("NotoKR", 9)
                py = 68
                for ln in wrap_text(c, prompt[:200], "NotoKR", 9, pw - 100):
                    c.drawString(50, py, ln)
                    py -= 12

            if data.get("tip"):
                c.setFillColor(colors.HexColor("#92400E"))
                c.setFont("NotoKR", 10)
                c.drawString(40, 18, f"💡 {data['tip']}")

        c.setFillColor(colors.HexColor("#94A3B8"))
        c.setFont("NotoKR", 10)
        c.drawRightString(pw - 25, 15, f"{i} / {len(SLIDES)}")
        c.showPage()

    c.save()
    print(f"PDF: {PDF_PATH}")


if __name__ == "__main__":
    build_pptx()
    build_pdf()
    print("Done.")
